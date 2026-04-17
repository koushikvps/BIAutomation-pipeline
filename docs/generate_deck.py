"""Generate the executive demo deck as a .pptx file in EY theme."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── EY Brand Colors ──
EY_YELLOW = RGBColor(0xFF, 0xE6, 0x00)
EY_BLACK = RGBColor(0x2E, 0x2E, 0x38)
EY_DARK = RGBColor(0x1A, 0x1A, 0x24)
EY_GRAY_DARK = RGBColor(0x33, 0x33, 0x40)
EY_GRAY_MID = RGBColor(0x74, 0x74, 0x80)
EY_GRAY_LIGHT = RGBColor(0xA0, 0xA0, 0xB0)
EY_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
EY_GREEN = RGBColor(0x00, 0xA7, 0x6D)
EY_RED = RGBColor(0xE0, 0x36, 0x2C)
EY_ORANGE = RGBColor(0xFF, 0x98, 0x31)
EY_BLUE = RGBColor(0x18, 0x8C, 0xE5)
EY_TEAL = RGBColor(0x00, 0xA3, 0xAE)
EY_PURPLE = RGBColor(0x7D, 0x3F, 0x98)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def add_bg(slide, color=EY_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_yellow_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(0.08), SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = EY_YELLOW
    shape.line.fill.background()


def add_ey_logo(slide, dark=False):
    left = Inches(11.5)
    top = Inches(0.3)
    txBox = slide.shapes.add_textbox(left, top, Inches(1.5), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "EY"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = EY_BLACK if dark else EY_YELLOW
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = "Building a better working world"
    run2.font.size = Pt(7)
    run2.font.color.rgb = RGBColor(0x50, 0x50, 0x50) if dark else EY_GRAY_MID


def add_section_label(slide, text, top=Inches(0.8), color=EY_YELLOW):
    txBox = slide.shapes.add_textbox(Inches(0.6), top, Inches(6), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.letter_spacing = Pt(3)


def add_title(slide, text, top=Inches(1.2), size=Pt(38), color=EY_WHITE, width=Inches(10)):
    txBox = slide.shapes.add_textbox(Inches(0.6), top, width, Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = True
    run.font.color.rgb = color
    return tf


def add_body(slide, text, top=Inches(2.5), left=Inches(0.6), width=Inches(11), size=Pt(16), color=EY_GRAY_LIGHT):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    return tf


def add_bullet_list(slide, items, top=Inches(2.5), left=Inches(0.6), width=Inches(5.5), size=Pt(14), color=EY_GRAY_LIGHT, bullet_color=None):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        # Determine prefix color
        if item.startswith("[YES]"):
            prefix = "\u2705 "
            item = item[5:].strip()
        elif item.startswith("[NO]"):
            prefix = "\u274C "
            item = item[4:].strip()
        elif item.startswith("[YELLOW]"):
            prefix = "\u25C9 "
            item = item[8:].strip()
        else:
            prefix = "\u25C9 "
        run = p.add_run()
        run.text = prefix + item
        run.font.size = size
        run.font.color.rgb = color
    return tf


def add_card(slide, title, body, left, top, width=Inches(3.8), height=Inches(3), accent_color=EY_YELLOW):
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = EY_GRAY_DARK
    shape.line.color.rgb = RGBColor(0x44, 0x44, 0x50)
    shape.line.width = Pt(1)
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    # Title
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = EY_WHITE
    # Body
    txBox2 = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.6), width - Inches(0.4), height - Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(body if isinstance(body, list) else body.split("\n")):
        p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p2.space_after = Pt(4)
        run2 = p2.add_run()
        run2.text = line
        run2.font.size = Pt(11)
        run2.font.color.rgb = EY_GRAY_LIGHT


def add_table(slide, headers, rows, top=Inches(2.8), left=Inches(0.6), width=Inches(12), row_height=Inches(0.38), font_size=Pt(10)):
    cols_count = len(headers)
    rows_count = len(rows) + 1
    col_width = int(width / cols_count)
    table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, Inches(row_height * rows_count / Inches(1)))
    table = table_shape.table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = EY_GRAY_DARK
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = EY_YELLOW

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = EY_DARK if r % 2 == 0 else RGBColor(0x22, 0x22, 0x2E)
            for p in cell.text_frame.paragraphs:
                p.font.size = font_size
                p.font.color.rgb = EY_GRAY_LIGHT

    return table_shape


def add_stat_box(slide, value, label, left, top, val_color=EY_YELLOW):
    txBox = slide.shapes.add_textbox(left, top, Inches(2), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(value)
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = val_color
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(11)
    run2.font.color.rgb = EY_GRAY_LIGHT


def add_flow_node(slide, text, desc, left, top, border_color=EY_GRAY_MID, width=Inches(1.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = EY_GRAY_DARK
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = EY_WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(8)
    run2.font.color.rgb = EY_GRAY_MID


def add_flow_arrow(slide, left, top):
    txBox = slide.shapes.add_textbox(left, top, Inches(0.3), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "\u2192"
    run.font.size = Pt(18)
    run.font.color.rgb = EY_YELLOW


def divider_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, EY_YELLOW)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(0.08), SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = EY_BLACK
    shape.line.fill.background()
    add_ey_logo(slide, dark=True)
    add_section_label(slide, "SECTION", color=RGBColor(0x50, 0x50, 0x50))
    add_title(slide, title, size=Pt(44), color=EY_BLACK)
    add_body(slide, subtitle, top=Inches(2.8), color=RGBColor(0x50, 0x50, 0x50), size=Pt(18))
    return slide


# ═══════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)

# Yellow accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(1.2), Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = EY_YELLOW
bar.line.fill.background()

add_section_label(slide, "Core Assurance \u00B7 Data & Analytics", top=Inches(1.0))
add_title(slide, "AI-Powered Data Platform\n& Test Automation", top=Inches(1.8), size=Pt(44))
add_body(slide, "13 autonomous AI agents. Two independent products. Self-healing. Self-maintaining.\nFrom business requirement to deployed data model and tested Power App.", top=Inches(3.6), size=Pt(18))

# Tech pills
pills = ["Azure Synapse", "Azure Data Factory", "Azure AI Foundry", "Playwright", "Azure DevOps", "Teams"]
txBox = slide.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(10), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
for pill in pills:
    run = p.add_run()
    run.text = f"  {pill}  "
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = EY_YELLOW

add_body(slide, "Assurance Datawarehouse  \u00B7  Executive Briefing  \u00B7  April 2026", top=Inches(6.5), size=Pt(12), color=EY_GRAY_MID)


# ═══════════════════════════════════════════════
# SLIDE 2: The Problem
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)
add_section_label(slide, "The Challenge")
add_title(slide, "Manual Data Engineering\nIs the Bottleneck")

add_card(slide, "BI Model Development", [
    "\u274C 2-3 weeks per medallion model",
    "\u274C Manual SQL writing, review, deployment",
    "\u274C Inconsistent naming and patterns",
    "\u274C No automated validation before deploy",
    "\u274C Knowledge siloed in individual engineers",
    "\u274C Backlog grows faster than capacity",
], Inches(0.6), Inches(3.0), Inches(5.8), Inches(3.2), EY_RED)

add_card(slide, "Power Apps & Data Testing", [
    "\u274C Test cases written manually over days",
    "\u274C Click-through testing every release",
    "\u274C No automated regression suites",
    "\u274C Data quality discovered after incidents",
    "\u274C No unified UI + data test framework",
    "\u274C Bug reporting to ADO is manual",
], Inches(6.8), Inches(3.0), Inches(5.8), Inches(3.2), EY_RED)

add_stat_box(slide, "2-3", "Weeks per Model", Inches(0.6), Inches(6.4), EY_RED)
add_stat_box(slide, "80%", "Manual Effort", Inches(3.2), Inches(6.4), EY_RED)
add_stat_box(slide, "0", "Automated Tests", Inches(5.8), Inches(6.4), EY_RED)


# ═══════════════════════════════════════════════
# SLIDE 3: The Solution
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)
add_section_label(slide, "The Solution")
add_title(slide, "One Story ID \u2192 Everything Automated")
add_body(slide, "Enter an Azure DevOps work item. 12 AI agents handle the rest.", top=Inches(2.2), size=Pt(16))

add_card(slide, "Product 1: BI Pipeline Automation", [
    "\u2705 Story to deployed medallion model in minutes",
    "\u2705 AI generates Bronze / Silver / Gold SQL",
    "\u2705 ADF incremental load pipelines auto-created",
    "\u2705 AI Code Review (7 categories) before deploy",
    "\u2705 Human review gate with approve / decline",
    "\u2705 Self-healing on deployment failures",
], Inches(0.6), Inches(3.0), Inches(5.8), Inches(3.5), EY_GREEN)

add_card(slide, "Product 2: Test Automation Platform", [
    "\u2705 AI routes stories to UI / Data / Both",
    "\u2705 Playwright drives browser on QA machine live",
    "\u2705 SQL data validation tests run server-side",
    "\u2705 Auto-creates ADO Test Plans, Cases, Bugs",
    "\u2705 Excel reports + Teams adaptive cards",
    "\u2705 Custom test categories managed in Web UI",
], Inches(6.8), Inches(3.0), Inches(5.8), Inches(3.5), EY_TEAL)

add_stat_box(slide, "13", "AI Agents", Inches(0.6), Inches(6.4))
add_stat_box(slide, "55", "HTTP Endpoints", Inches(3.0), Inches(6.4))
add_stat_box(slide, "2", "Products", Inches(5.4), Inches(6.4))
add_stat_box(slide, "1 Day", "To Integrate", Inches(8.0), Inches(6.4), EY_GREEN)


# ═══════════════════════════════════════════════
# SLIDE 4: Impact Numbers
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)
add_section_label(slide, "Business Impact")
add_title(slide, "Measurable Acceleration")

add_stat_box(slide, "95%", "Reduction in Build Time", Inches(0.6), Inches(2.0), EY_GREEN)
add_stat_box(slide, "10x", "Story Throughput", Inches(4.0), Inches(2.0), EY_GREEN)
add_stat_box(slide, "100%", "Automated Validation", Inches(7.5), Inches(2.0), EY_GREEN)

add_table(slide, ["Metric", "Before (Manual)", "After (AI Platform)", "Improvement"], [
    ["Time per medallion model", "2-3 weeks", "10-15 minutes", "95% faster"],
    ["Stories per sprint per engineer", "1-2", "10-20", "10x throughput"],
    ["Code review coverage", "Varies", "100% (7 categories)", "Complete"],
    ["Data quality validation", "Post-incident", "Every deployment", "Proactive"],
    ["Naming consistency", "Team-dependent", "Convention-enforced", "100%"],
    ["Test creation time", "Days per story", "Minutes (AI-generated)", "95% faster"],
    ["Regression testing", "Manual per release", "Automated, replayable", "Continuous"],
], top=Inches(3.5))


# ═══════════════════════════════════════════════
# SLIDE 5: Divider - Architecture
# ═══════════════════════════════════════════════
divider_slide("Platform Architecture", "Two independent products. One shared infrastructure. Zero vendor lock-in.")


# ═══════════════════════════════════════════════
# SLIDE 6: Two-Product Architecture
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)
add_section_label(slide, "Architecture")
add_title(slide, "Two Independent Products")
add_body(slide, "Deploy, scale, and manage each product independently. Share infrastructure costs.", top=Inches(2.2), size=Pt(15))

add_card(slide, "Product 1: BI Pipeline", [
    "Function App: {prefix}-{env}-func",
    "44 HTTP Endpoints + 3 Timer Triggers",
    "2 Orchestrators (story: 9 steps, fix: 8 steps)",
    "7 AI Agents (incl. Discovery + Bug Fixer)",
    "Ops Module: auto-pause, secret alerts, cleanup",
    "Durable Task Hub: BiAutoHubV8",
], Inches(0.6), Inches(3.0), Inches(5.8), Inches(3.2), EY_YELLOW)

add_card(slide, "Product 2: Test Automation", [
    "Function App: {prefix}-{env}-test-func",
    "11 HTTP Endpoints",
    "1 Orchestrator (6 steps)",
    "6 AI Agents",
    "Durable Task Hub: BiAutoTestHubV1",
], Inches(6.8), Inches(3.0), Inches(5.8), Inches(3.2), EY_TEAL)

# Shared infra line
add_body(slide, "Shared: App Service Plan (EP1)  \u00B7  Storage Account  \u00B7  Key Vault  \u00B7  App Insights  \u00B7  Azure SQL  \u00B7  Synapse  \u00B7  Data Factory", top=Inches(6.4), size=Pt(12), color=EY_YELLOW)


# ═══════════════════════════════════════════════
# SLIDE 7: BI Pipeline Flow
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)
add_section_label(slide, "Product 1")
add_title(slide, "BI Pipeline \u2014 Agent Flow")
add_body(slide, "9-step story pipeline + 8-step bug fix pipeline. Self-correcting loop.", top=Inches(2.2), size=Pt(15))

nodes_row1 = [
    ("ADO Story", "Fetch work item", EY_GRAY_MID),
    ("Planner", "Analyze + plan", EY_YELLOW),
    ("Review Gate", "Human approve", EY_ORANGE),
    ("Developer", "Generate SQL + ADF", EY_YELLOW),
    ("Code Review", "7-category AI check", EY_YELLOW),
]
x = Inches(0.4)
for name, desc, color in nodes_row1:
    add_flow_node(slide, name, desc, x, Inches(3.0), color)
    x += Inches(1.7)
    if name != "Code Review":
        add_flow_arrow(slide, x - Inches(0.2), Inches(3.0))

nodes_row2 = [
    ("Pre-Validate", "Syntax check", EY_GRAY_MID),
    ("Deploy", "Synapse + ADF", EY_GREEN),
    ("Post-Validate", "DQ scoring", EY_GREEN),
    ("Healer", "Self-fix failures", EY_YELLOW),
]
x = Inches(2.0)
for name, desc, color in nodes_row2:
    add_flow_node(slide, name, desc, x, Inches(4.2), color)
    x += Inches(1.7)
    if name != "Healer":
        add_flow_arrow(slide, x - Inches(0.2), Inches(4.2))

# Bug Fix Loop
add_body(slide, "\u21BB  Self-Correcting Loop:  Bug Found \u2192 Bug Fixer Agent \u2192 AI Fix \u2192 Code Review \u2192 Re-Deploy \u2192 Re-Test \u2192 Resolved", top=Inches(5.4), size=Pt(13), color=EY_ORANGE)

# Quote
txBox = slide.shapes.add_textbox(Inches(0.6), Inches(6.0), Inches(11), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "\u201CIf a bug is found, the Bug Fixer Agent reads the ADO bug, generates corrected code, and re-deploys \u2014 no human coding required.\u201D"
run.font.size = Pt(13)
run.font.italic = True
run.font.color.rgb = EY_GRAY_LIGHT


# ═══════════════════════════════════════════════
# SLIDE 8: Test Automation Flow
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)
add_section_label(slide, "Product 2")
add_title(slide, "Test Automation \u2014 Agent Flow")
add_body(slide, "AI-powered testing for Power Apps (UI) and data warehouse (SQL). Routes intelligently.", top=Inches(2.2), size=Pt(15))

test_nodes = [
    ("Story", "Any format", EY_GRAY_MID),
    ("Router", "UI / Data / Both", EY_YELLOW),
    ("Test Planner", "Gen scenarios", EY_YELLOW),
    ("Generator", "Playwright scripts", EY_YELLOW),
    ("Executor", "Run on QA machine", EY_GREEN),
    ("Reporter", "ADO + Excel + Teams", EY_GREEN),
]
x = Inches(0.3)
for name, desc, color in test_nodes:
    add_flow_node(slide, name, desc, x, Inches(3.2), color, Inches(1.8))
    x += Inches(2.0)
    if name != "Reporter":
        add_flow_arrow(slide, x - Inches(0.2), Inches(3.2))

add_card(slide, "UI Testing", [
    "QA watches agent drive Edge in",
    "real-time. Headed mode, visible",
    "clicks, live streaming to Web UI.",
], Inches(0.6), Inches(4.6), Inches(3.6), Inches(2.0), EY_TEAL)

add_card(slide, "Data Testing", [
    "SQL validation: row counts, nulls,",
    "duplicates, referential integrity,",
    "cross-layer consistency. Server-side.",
], Inches(4.6), Inches(4.6), Inches(3.6), Inches(2.0), EY_BLUE)

add_card(slide, "Reporting", [
    "Auto-creates ADO Test Plans +",
    "Cases + Runs. Files bugs. Excel",
    "workbook + Teams adaptive cards.",
], Inches(8.6), Inches(4.6), Inches(3.6), Inches(2.0), EY_PURPLE)


# ═══════════════════════════════════════════════
# SLIDE 9: Divider - Integration Mode
# ═══════════════════════════════════════════════
divider_slide("Integration Mode", "Plug into the platform you already have. No rip-and-replace. No migration.")


# ═══════════════════════════════════════════════
# SLIDE 10: Integration Mode - Why
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)
add_section_label(slide, "Integration Mode")
add_title(slide, "Your Platform. Our Intelligence.")
add_body(slide, "Enterprise teams already have mature data platforms. They don't need another one. They need acceleration.", top=Inches(2.2), size=Pt(16))

add_card(slide, "What Competitors Ask You to Do", [
    "\u274C Migrate to their platform (Fabric, Informatica)",
    "\u274C Learn their proprietary interface",
    "\u274C Replace your CI/CD with theirs",
    "\u274C Abandon your naming conventions",
    "\u274C Re-train your engineering team",
    "\u274C 6-12 month implementation timeline",
], Inches(0.6), Inches(3.2), Inches(5.8), Inches(3.5), EY_RED)

add_card(slide, "What We Do", [
    "\u2705 Connect to your existing Synapse, ADF, ADLS",
    "\u2705 Auto-detect your naming conventions",
    "\u2705 Generate code that matches your style",
    "\u2705 Deliver through your CI/CD as a PR",
    "\u2705 Your engineers review, your pipeline deploys",
    "\u2705 Day 1 integration",
], Inches(6.8), Inches(3.2), Inches(5.8), Inches(3.5), EY_GREEN)


# ═══════════════════════════════════════════════
# SLIDE 11: Integration Mode - How
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)
add_section_label(slide, "Integration Mode")
add_title(slide, "Four-Step Onboarding")

cards = [
    ("Step 1: Discover", "Discovery Agent scans your Synapse (schemas, tables, views, procs, distributions), ADF (pipelines, datasets, triggers), and ADLS (folders, formats). 10-15 min, fully automated.", EY_YELLOW),
    ("Step 2: Learn Conventions", "Convention Adapter auto-detects naming patterns: prefixes, case, schema purposes, distributions. Rules generated, you can override.", EY_YELLOW),
    ("Step 3: Submit Stories", "Same workflow. Planner and Developer now use YOUR catalog and YOUR conventions. Generated SQL looks like your team wrote it.", EY_YELLOW),
    ("Step 4: Receive PR", "Artifacts delivered as Pull Request to your ADO repo. Your team reviews. Your CI/CD deploys. Your governance stays intact.", EY_GREEN),
]
x = Inches(0.4)
for title, body, color in cards:
    add_card(slide, title, [body], x, Inches(2.5), Inches(3.0), Inches(3.2), color)
    x += Inches(3.15)

# Quote
txBox = slide.shapes.add_textbox(Inches(0.6), Inches(6.0), Inches(11), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "\u201CIt already knows our naming conventions, our folder structure, our pipeline patterns. The SQL it generated looks like our team wrote it.\u201D"
run.font.size = Pt(14)
run.font.italic = True
run.font.color.rgb = EY_GRAY_LIGHT


# ═══════════════════════════════════════════════
# SLIDE 12: Convention Detection
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)
add_section_label(slide, "Integration Mode")
add_title(slide, "Auto-Detected Conventions")
add_body(slide, "The Discovery Agent reads your environment and generates rules. All agents follow them.", top=Inches(2.2), size=Pt(15))

add_table(slide, ["Detection", "Example", "Effect on Generated Code"], [
    ["Naming case", "raw_customer_orders detected", "All objects use snake_case"],
    ["Table prefix", "80% of tables start with tbl_", "Generated tables get tbl_ prefix"],
    ["View prefix", "90% of views start with vw_", "Generated views get vw_ prefix"],
    ["Bronze schema", "Schema [raw] has external tables", "Bronze layer targets [raw]"],
    ["Silver schema", "Schema [cleansed] has tables + procs", "Silver layer targets [cleansed]"],
    ["Gold schema", "Schema [analytics] is view-heavy", "Gold layer targets [analytics]"],
    ["Distribution", "70% use HASH distribution", "Default distribution set to HASH"],
    ["ADF naming", "Pipelines follow pl_{src}_{tgt}", "Generated pipelines match pattern"],
    ["Delivery mode", "Client has CI/CD pipeline", "Artifacts delivered as PR"],
], top=Inches(2.8))


# ═══════════════════════════════════════════════
# SLIDE 13: Three Modes
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)
add_section_label(slide, "Platform Modes")
add_title(slide, "Three Modes, One Platform")

add_table(slide, ["Aspect", "Greenfield", "Brownfield", "Integration Mode"], [
    ["Target", "New platform", "Our existing deploy", "Client's existing platform"],
    ["Discovery", "N/A", "Check existing tables", "Full environment scan"],
    ["Conventions", "Our defaults", "Our conventions", "Client's auto-detected"],
    ["Delivery", "Direct deploy", "Direct deploy", "Pull Request to client repo"],
    ["Schema names", "bronze/silver/gold", "bronze/silver/gold", "Client's actual names"],
    ["CI/CD", "Our pipeline", "Our pipeline", "Client's existing CI/CD"],
    ["Risk", "None (new)", "Low (additive)", "Minimal (PR review gate)"],
    ["Onboarding", "Deploy infra", "Deploy code", "Connect + Discover (1 day)"],
], top=Inches(2.4))


# ═══════════════════════════════════════════════
# SLIDE 14: Divider - Competition
# ═══════════════════════════════════════════════
divider_slide("Competitive Landscape", "How we compare to every alternative on the market.")


# ═══════════════════════════════════════════════
# SLIDE 15: Competitive Matrix
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)
add_section_label(slide, "Competitive Analysis")
add_title(slide, "Feature Comparison")

Y = "\u25CF"  # filled = yes
N = "\u25CB"  # empty = no
P = "\u25D0"  # half = partial

add_table(slide, ["Capability", "Our Platform", "Fabric Copilot", "Informatica", "dbt", "Coalesce", "Matillion", "Prophecy"], [
    ["Story-to-SQL automation",      Y+" Full",  N+" No",  N+" No",  N+" No",  N+" No",  N+" No",  P+" Partial"],
    ["Medallion architecture gen",   Y+" Full",  N+" No",  N+" No",  P+" Manual", P+" Manual", P+" Manual", P+" Manual"],
    ["Existing environment scan",    Y+" Full",  N+" No",  Y+" Catalog", N+" No", N+" No",  N+" No",  N+" No"],
    ["Auto-detect conventions",      Y+" Full",  N+" No",  N+" No",  N+" No",  N+" No",  N+" No",  N+" No"],
    ["Multi-agent SDLC",             Y+" 13 agents", N+" No", P+" Specialized", N+" No", N+" No", N+" No", N+" No"],
    ["Self-healing pipeline",        Y+" Yes",   N+" No",  N+" No",  N+" No",  N+" No",  N+" No",  N+" No"],
    ["PR delivery mode",             Y+" Yes",   N+" No",  N+" No",  Y+" CI",  Y+" Git", N+" No",  Y+" Git"],
    ["Azure Synapse native",         Y+" Yes",   P+" Fabric", P+" Connector", P+" Adapter", N+" No", Y+" Yes", N+" No"],
    ["Test automation (UI+Data)",    Y+" Full",  N+" No",  N+" No",  P+" Data", N+" No",  N+" No",  N+" No"],
    ["Works with existing stack",    Y+" Day 1", N+" Migrate", P+" 6-12mo", P+" Weeks", P+" Weeks", P+" Weeks", P+" Weeks"],
], top=Inches(2.2), font_size=Pt(9))


# ═══════════════════════════════════════════════
# SLIDE 16: Why Not Competitors
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)
add_section_label(slide, "Competitive Analysis")
add_title(slide, "Why Not the Alternatives?")

add_card(slide, "Microsoft Fabric Copilot", [
    "Smart autocomplete, not autonomous.",
    "Requires F64 capacity ($5K+/mo).",
    "Must migrate from Synapse to Fabric.",
    "",
    "A copilot helps you type. We do the typing.",
], Inches(0.6), Inches(2.5), Inches(5.8), Inches(2.5), EY_RED)

add_card(slide, "Informatica CLAIRE", [
    "Governance platform, not warehouse automation.",
    "$200K-$500K+ Year 1. 6-12 month implementation.",
    "CLAIRE Agents don't write Synapse DDL.",
    "",
    "A governance tool. We're a dev accelerator.",
], Inches(6.8), Inches(2.5), Inches(5.8), Inches(2.5), EY_RED)

add_card(slide, "dbt / Coalesce / Matillion", [
    "Frameworks for organizing manual work.",
    "You still write every SQL model / drag every node.",
    "No AI generation from business requirements.",
    "",
    "Better tooling for engineers. We eliminate the work.",
], Inches(0.6), Inches(5.2), Inches(5.8), Inches(2.0), EY_RED)

add_card(slide, "Build It Yourself", [
    "3-5 engineers, 6-12 months, $500K-$1M.",
    "LLM call is 10%. Orchestration + Synapse",
    "quirks + self-healing is the other 90%.",
    "",
    "We've already built it. 13 agents. 55 endpoints.",
], Inches(6.8), Inches(5.2), Inches(5.8), Inches(2.0), EY_RED)


# ═══════════════════════════════════════════════
# SLIDE 17: Cost Comparison
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)
add_section_label(slide, "Economics")
add_title(slide, "Cost Comparison")

add_table(slide, ["Solution", "Year 1 Cost", "Annual Ongoing", "Time to Value", "Migration Required?"], [
    ["\u2B50 Our Platform", "$500-1,500/mo Azure", "Same", "1 Day", "No"],
    ["Fabric Copilot", "$60K+ (F64 capacity)", "$60K+", "3-6 months", "Yes (to Fabric)"],
    ["Informatica CLAIRE", "$200K-500K+", "$150K+", "6-12 months", "Yes (IDMC platform)"],
    ["dbt Cloud Enterprise", "$50K+", "$50K+", "1-3 months", "New runtime"],
    ["Coalesce Enterprise", "$30K+", "$30K+", "1-3 months", "New platform"],
    ["Matillion Enterprise", "$40K-100K+", "$40K+", "1-3 months", "New runtime"],
    ["Internal Build", "$500K-1M (engineers)", "$200K+ maintenance", "6-12 months", "N/A"],
], top=Inches(2.4))

# Quote
txBox = slide.shapes.add_textbox(Inches(0.6), Inches(5.8), Inches(11), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "\u201CThe only solution that costs nothing beyond Azure infrastructure, requires no migration, and delivers value on Day 1.\u201D"
run.font.size = Pt(14)
run.font.italic = True
run.font.color.rgb = EY_YELLOW


# ═══════════════════════════════════════════════
# SLIDE 18: Divider - Technology
# ═══════════════════════════════════════════════
divider_slide("Technology Stack", "Built entirely on Azure. No third-party platforms. No vendor lock-in.")


# ═══════════════════════════════════════════════
# SLIDE 19: Azure Resources
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)
add_section_label(slide, "Infrastructure")
add_title(slide, "Azure Resource Inventory")
add_body(slide, "13 resource types. All deployed via Bicep IaC. Single az deployment command.", top=Inches(2.2), size=Pt(15))

resources = [
    ("Synapse Dedicated Pool", "DW100c \u00B7 Medallion DW"),
    ("Function App (BI)", "Premium EP1 \u00B7 44 endpoints + 3 timers"),
    ("Function App (Test)", "Shared EP1 \u00B7 11 endpoints"),
    ("Azure SQL Database", "Config DB + Source Data"),
    ("Key Vault", "Secrets + RBAC for all identities"),
    ("ADLS Gen2", "Bronze data lake + artifacts"),
    ("Data Factory", "Incremental load pipelines"),
    ("AI Foundry", "Phi-4 model \u00B7 LLM for all agents"),
    ("App Insights", "Telemetry + monitoring"),
    ("Log Analytics", "90-day centralized logs"),
]
col1 = resources[:5]
col2 = resources[5:]

for i, (name, desc) in enumerate(col1):
    add_card(slide, name, [desc], Inches(0.5), Inches(2.8 + i * 0.85), Inches(5.8), Inches(0.7), EY_YELLOW)
for i, (name, desc) in enumerate(col2):
    add_card(slide, name, [desc], Inches(6.8), Inches(2.8 + i * 0.85), Inches(5.8), Inches(0.7), EY_YELLOW)


# ═══════════════════════════════════════════════
# SLIDE 20: Security
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)
add_section_label(slide, "Enterprise Readiness")
add_title(slide, "Security & Governance")

sec_cards = [
    ("Zero Secrets in Code", "All credentials in Key Vault. Managed Identity for service-to-service. No passwords in config or logs.", EY_YELLOW),
    ("SQL Injection Prevention", "All dynamic SQL sanitized. Parameterized queries. Input validation at every endpoint.", EY_YELLOW),
    ("Human Review Gate", "AI generates, humans approve. Every pipeline pauses for review before deployment.", EY_GREEN),
    ("Full Audit Trail", "Every execution logged: timestamp, user, story ID, artifacts, validation scores, deploy status.", EY_YELLOW),
    ("RBAC + TLS 1.2", "Managed identities with least-privilege. HTTPS-only. FTPS disabled. Function key auth.", EY_YELLOW),
    ("PR Delivery for Compliance", "Integration Mode delivers through client's change management. PR = audit trail + sign-off.", EY_GREEN),
]
positions = [
    (Inches(0.5), Inches(2.3)), (Inches(4.4), Inches(2.3)), (Inches(8.3), Inches(2.3)),
    (Inches(0.5), Inches(4.6)), (Inches(4.4), Inches(4.6)), (Inches(8.3), Inches(4.6)),
]
for (name, desc, color), (x, y) in zip(sec_cards, positions):
    add_card(slide, name, [desc], x, y, Inches(3.6), Inches(2.0), color)


# ═══════════════════════════════════════════════
# SLIDE 21: Roadmap
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)
add_section_label(slide, "What's Next")
add_title(slide, "Product Roadmap")

add_card(slide, "Delivered (POC \u2014 v6.0)", [
    "\u2705 13 AI agents across 2 products",
    "\u2705 Greenfield + Brownfield + Integration Mode",
    "\u2705 Discovery Agent + Convention Adapter",
    "\u2705 Self-correcting Bug Fixer loop",
    "\u2705 Ops Module (auto-pause, secrets, cleanup)",
    "\u2705 Human review gate + Cancel flow",
    "\u2705 Teams adaptive card notifications",
    "\u2705 UI + Data test automation",
    "\u2705 Web UI with 16 pages (incl. Ops Dashboard)",
    "\u2705 127 unit tests, full Bicep IaC",
], Inches(0.5), Inches(2.3), Inches(3.8), Inches(4.8), EY_GREEN)

add_card(slide, "Phase 2: Production", [
    "\u25C9 Multi-tenant SaaS deployment",
    "\u25C9 Power BI report auto-generation",
    "\u25C9 Databricks / Snowflake support",
    "\u25C9 Advanced semantic layer (metrics, KPIs)",
    "\u25C9 Git-based artifact version control",
    "\u25C9 Role-based Web UI access control",
    "\u25C9 Custom LLM fine-tuning on client data",
], Inches(4.7), Inches(2.3), Inches(3.8), Inches(4.8), EY_TEAL)

add_card(slide, "Phase 3: Scale", [
    "\u25C9 Multi-cloud (AWS Redshift, GCP BigQuery)",
    "\u25C9 Azure Marketplace listing",
    "\u25C9 Self-service onboarding portal",
    "\u25C9 AI-driven performance optimization",
    "\u25C9 Real-time streaming pipeline support",
    "\u25C9 Industry-specific templates",
    "\u25C9 EY Consulting marketplace listing",
], Inches(8.9), Inches(2.3), Inches(3.8), Inches(4.8), EY_PURPLE)


# ═══════════════════════════════════════════════
# SLIDE 22: Closing
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_yellow_bar(slide)
add_ey_logo(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(1.2), Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = EY_YELLOW
bar.line.fill.background()

add_title(slide, "Every other tool helps your\nengineers work faster.", top=Inches(1.8), size=Pt(36))

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(3.2), Inches(10), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "We do the work."
run.font.size = Pt(40)
run.font.bold = True
run.font.color.rgb = EY_YELLOW

add_stat_box(slide, "13", "AI Agents", Inches(0.5), Inches(4.4))
add_stat_box(slide, "55", "Endpoints", Inches(2.7), Inches(4.4))
add_stat_box(slide, "3", "Deploy Modes", Inches(5.0), Inches(4.4))
add_stat_box(slide, "1 Day", "To Integrate", Inches(7.3), Inches(4.4), EY_GREEN)
add_stat_box(slide, "$0", "New Licenses", Inches(9.7), Inches(4.4), EY_GREEN)

# Quote
txBox = slide.shapes.add_textbox(Inches(0.6), Inches(5.8), Inches(11), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "\u201CYour platform. Your conventions. Your CI/CD. Our intelligence.\u201D"
run.font.size = Pt(20)
run.font.bold = True
run.font.color.rgb = EY_GRAY_LIGHT

add_body(slide, "EY Core Assurance  \u00B7  Data & Analytics  \u00B7  AI Automation Practice\nReady for live demo", top=Inches(6.5), size=Pt(13), color=EY_GRAY_MID)


# ═══════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "demo-deck.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
